from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR,PP_ALIGN
width=11.02
height=6.2

# Assuming this is the LangChain output
langchain_output = [{'title': 'Leapfrog Anti-Harassment Policy (Nepal)'}, {'header': 'Introduction', 'points': ['Dedicated to safe and respectful workplace', 'Zero tolerance for harassment or discrimination', 'Outline of harassment types and reporting steps']}, {'header': 'Purpose', 'paragraph': 'The purpose of this policy is to protect individuals from various forms of harassment, ensure fair handling of incidents, and establish a clear process for reporting and addressing complaints.'}, {'header': 'Scope', 'points': ['Applies to all employees and contractors', 'Covers work-related activities beyond the office', 'Includes remote work and company-sponsored events']}, {'header': 'Definitions', 'paragraph': 'This section defines harassment, discrimination, and sexual harassment to establish a common understanding of unacceptable behavior.'}, {'header': 'Means of Harassment', 'points': ['Verbal, Non-Verbal, and Physical Harassment', 'Cyber Harassment via online platforms', 'Unwanted behaviors that intimidate or demean individuals']}, {'header': 'Reporting Mechanism', 'paragraph': 'Employees experiencing harassment can report incidents through various channels such as the Harassment Reporting form and direct reporting to superiors, ensuring confidentiality and appropriate action.'}, {'header': 'Investigation Mechanism', 'points': ['Involves formal investigation by the Board of Directors', 'External legal consultants appointed for sexual harassment cases', 'Completion of investigation within specified timelines']}, {'header': 'Confidentiality', 'paragraph': 'All complaints and investigations are kept strictly confidential, protecting the identities of complainants, accused, and witnesses.'}, {'header': 'Protection of the Complainant', 'points': ['No adverse action against complainants', 'Temporary arrangements for safety during investigations']}, {'header': 'Retaliation', 'paragraph': 'Any form of retaliation against complainants or witnesses is prohibited, with clear reporting mechanisms and prompt investigations into such claims.'}, {'header': 'Response and Resolution', 'points': ['Corrective actions based on severity of offenses', 'Resolution process to be completed within 14 days']}, {'header': 'Conciliation', 'paragraph': 'Parties can voluntarily resolve harassment complaints through conciliation efforts, with confidential agreements in place upon successful resolution.'}, {'header': 'Expectations from Employees', 'points': ['Professional conduct expected at all times', 'Bystanders must report observed harassment']}, {'header': 'Prevention and Awareness', 'paragraph': 'Regular training and communication about the policy will be conducted to promote awareness and a respectful workplace culture.'}, {'header': 'Disciplinary Action', 'points': ['Consequences for offenders may include termination', 'Protection against action for false complaints if not maliciously intended']}, {'header': 'Legal Framework', 'paragraph': 'This policy aligns with Nepal’s Labor Act and pertinent legal provisions, ensuring compliance and protection of rights.'}, {'header': 'Revision History', 'points': ['Policy revisions are recorded with responsible personnel.', 'Annual reviews planned to ensure alignment with legal changes.']}, {'header': 'Conclusion', 'paragraph': 'The Leapfrog Anti-Harassment Policy aims to foster a safe work environment by strictly prohibiting harassment and outlining clear reporting, investigation, and resolution mechanisms.'}, {'title': '“A safe workplace is a right for every employee.”'}]

def generate_ppt(langchain_output,filename):
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    for section in langchain_output:
        # Add new slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        if "title" in section:
            # Title slide (first item, assumed)
            title = section["title"]
            
            # Add green rectangle for background
            header_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=Inches(0),
                top=Inches(0),
                width=Inches(width),
                height=Inches(height)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = RGBColor(0, 153, 51)  # Green color

            # Add centered title text
            header_textbox = slide.shapes.add_textbox(
                left=Inches(0),
                top=Inches(height/2-1),  # Adjust this to better center the text
                width=Inches(11.02),
                height=Inches(height)
            )
            header_text = header_textbox.text_frame.add_paragraph()
            header_text.text = title
            header_text.font.color.rgb = RGBColor(255, 255, 255)  # White color
            header_text.font.size = Pt(32)
            header_text.font.bold = True
            tf = header_textbox.text_frame
            tf.word_wrap = True
            header_text.alignment = PP_ALIGN.CENTER  # Center the text horizontally


        if 'header' in section:
        # Title slide (first item, assumed)
            title = section["header"]

            # Add green rectangle for header
            header_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=Inches(0),
                top=Inches(0),
                width=Inches(11.02),
                height=Inches(height)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = RGBColor(0, 153, 51)  # Green color

            # Add header text
            header_textbox = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(-0.05),
                width=Inches(width),
                height=Inches(0.5)
            )
            tf = header_textbox.text_frame  # Access the text frame
            tf.word_wrap = True             # Enable word wrapping for the text frame
            tf.horizontal_anchor = PP_ALIGN.LEFT  # Align the entire text frame to the left
            header_text = tf.add_paragraph()
            header_text.text = title
            header_text.font.color.rgb = RGBColor(255, 255, 255)  # White color
            header_text.font.size = Pt(25)
            header_text.font.bold = True
            header_text.alignment = PP_ALIGN.LEFT  # Align the specific paragraph to the left


        if "points" in section:
            # Add section with header and bullet points
            points = section["points"]
            # Add green rectangle for content (rounded rectangle)
            content_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=Inches(-1),
                top=Inches(1),
                width=Inches(10.8),
                height=Inches(height)
            )
            content_shape.fill.solid()
            content_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            content_shape.line.fill.background()  # Remove the border
            content_shape.adjustments[0] = 0.1  # Round top-right corner
            
            # Add content textbox
            content_textbox = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(1.5),
                width=Inches(8),
                height=Inches(5)
            )
            tf = content_textbox.text_frame
            tf.word_wrap = True
            # Add bullet points
            for index,point in enumerate(points):
                p = tf.add_paragraph()
                p.text = str(index+1)+" "+point
                p.level = 1  # Bullet point level
                p.font.size = Pt(14)
                p.space_after = Pt(12)

        if "paragraph" in section:
            # Add section with header and paragraph (no bullet points)
            paragraph = section["paragraph"]
            
            # Add green rectangle for content (rounded rectangle)
            content_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=Inches(-1),
                top=Inches(1),
             width=Inches(10.8),
                height=Inches(height)
            )
            content_shape.fill.solid()
            content_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            content_shape.line.fill.background()  # Remove the border
            content_shape.adjustments[0] = 0.1  # Round top-right corner
            
            # Add content textbox
            content_textbox = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(1.5),
                width=Inches(8),
                height=Inches(5)
            )
            tf = content_textbox.text_frame
            tf.word_wrap = True

            # Add paragraph content (no bullet points)
            p = tf.add_paragraph()
            p.text = paragraph
            p.font.size = Pt(14)
            p.space_after = Pt(12)

    # Save the presentation
    prs.save(f"{filename}")

# Call the function with LangChain output
if __name__=="__main__":
    generate_ppt(langchain_output,"slides.pptx")
